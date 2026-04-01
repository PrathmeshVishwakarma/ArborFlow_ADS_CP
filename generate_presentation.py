#!/usr/bin/env python3
"""
ArborFlow PowerPoint Presentation Generator
Generates comprehensive presentation about ArborFlow project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Color scheme
TITLE_COLOR = RGBColor(25, 118, 210)
ACCENT_COLOR = RGBColor(244, 67, 54)
TEXT_COLOR = RGBColor(66, 66, 66)
BG_COLOR = RGBColor(245, 245, 245)

def add_title_slide(title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, bullet_points):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Separator line
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(1.1), Inches(9.5), Inches(1.1))
    line.line.color.rgb = ACCENT_COLOR
    line.line.width = Pt(2)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        if isinstance(point, tuple):
            p.text = point[0]
            p.level = point[1] if len(point) > 1 else 0
        else:
            p.text = point
            p.level = 0
        
        p.font.size = Pt(20) if p.level == 0 else Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(4)
        p.space_after = Pt(4)

def add_two_column_slide(title, left_title, left_points, right_title, right_points):
    """Add slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    for point in left_points:
        p = left_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_after = Pt(4)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    for point in right_points:
        p = right_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_after = Pt(4)

# PRESENTATION SLIDES

# Slide 1: Title
add_title_slide(
    "ArborFlow",
    "High-Performance Network Processing Engine\nBuilt with Advanced Data Structures"
)

# Slide 2: Project Objectives
add_content_slide(
    "Project Objectives",
    [
        "Design and implement a real-time network security engine in C",
        "Capture and filter network packets at millions per second",
        ("Implement efficient IP filtering using advanced data structures:", 0),
        ("  - Bit Vectors: O(1) prefix filtering", 1),
        ("  - vEB Trees: O(1) amortized exact IP matching", 1),
        "Process packets with sub-microsecond latency per packet",
        "Demonstrate integration of networking, algorithms, and systems programming",
        "Create production-ready, modular firewall components"
    ]
)

# Slide 3: Architecture Overview
add_content_slide(
    "Architecture Overview",
    [
        ("Pipeline: Internet Traffic → Output Processing", 0),
        ("1. Capture Layer (libpcap + pthreads)", 0),
        ("   - Intercept raw packets from network interface", 1),
        ("   - Parse: Ethernet → IP → TCP/UDP headers", 1),
        ("2. Queue Layer (Lock-Free SPSC)", 0),
        ("   - Thread-safe packet passing with atomics", 1),
        ("3. Filtering Layer (Gatekeeper)", 0),
        ("   - Layer 1: BitVector prefix check O(1)", 1),
        ("   - Layer 2: vEB Tree exact match O(1)", 1),
        ("4. Scheduling Layer (Max Heap)", 0),
        ("   - Priority-based packet ordering for QoS", 1)
    ]
)

# Slide 4: Tech Stack
add_two_column_slide(
    "Technology Stack",
    "Core Engine",
    [
        "C (C11 Standard)",
        "libpcap: Packet capture",
        "pthreads: Multi-threading",
        "C11 stdatomic.h: Lock-free ops",
        "GCC -O2: Optimized compilation"
    ],
    "ML Integration (Optional)",
    [
        "Python 3.10+",
        "XGBoost: Attack classification",
        "Scikit-learn/Isolation Forest",
        "Pandas: Data processing",
        "IPC: Subprocess communication"
    ]
)

# Slide 5: Methodology
add_content_slide(
    "Methodology",
    [
        ("Phase 1: Literature Review", 0),
        ("  - vEB Trees, Bit Vectors, Lock-free algorithms", 1),
        ("Phase 2: Architecture Design", 0),
        ("  - Hierarchical modular system design", 1),
        ("Phase 3: Implementation", 0),
        ("  - Component development: Capture → Queue → Filter → Schedule", 1),
        ("Phase 4: Testing", 0),
        ("  - Unit tests, stress tests (1M packets), integration tests", 1),
        ("Phase 5: Optimization", 0),
        ("  - Memory efficiency, lock-free concurrency", 1)
    ]
)

# Slide 6: Literature Review
add_content_slide(
    "Literature Review",
    [
        ("Bit Vectors (Knuth, 1969)", 0),
        ("  Dense boolean arrays with O(1) membership testing", 1),
        ("Van Emde Boas Trees (van Emde Boas, 1975)", 0),
        ("  O(log log U) operations; two-tier variant for IP space", 1),
        ("Lock-Free Data Structures (Herlihy & Shavit, 2008)", 0),
        ("  Concurrent programming without mutexes using atomics", 1),
        ("Network Protocol Processing (Stevens, 1994)", 0),
        ("  Packet layering, encapsulation, header parsing", 1),
        ("Firewall & IDS Design (Cisco, Palo Alto)", 0),
        ("  Multi-layer defense, real-time filtering, threat detection", 1)
    ]
)

# Slide 7: Packet Collection (Member 1)
add_content_slide(
    "Member 1: Packet Collection Engine",
    [
        ("Technology: libpcap + pthreads", 0),
        ("Functionality:", 0),
        ("  1. Open network interface in promiscuous mode", 1),
        ("  2. Apply BPF (Berkeley Packet Filter) to kernel", 1),
        ("  3. Capture raw Ethernet frames at line rate", 1),
        ("  4. Parse headers: Ethernet (14B) → IP (20B) → TCP/UDP (20/8B)", 1),
        ("Output per packet: src_ip, dst_ip, protocol, size, priority", 0),
        ("Performance: No packet loss in normal conditions", 0)
    ]
)

# Slide 8: Concurrent Queue (Member 1/2)
add_content_slide(
    "Lock-Free Concurrent Queue",
    [
        ("Design: SPSC (Single-Producer, Single-Consumer)", 0),
        ("Capacity: 1024 packets (circular buffer)", 1),
        ("Synchronization: C11 atomic_uint (no mutexes)", 1),
        ("Why Lock-Free?", 0),
        ("  - Mutexes cause context switches → 1000s of nanoseconds", 1),
        ("  - Atomics are CPU instructions → 10-100 nanoseconds", 1),
        ("Operations:", 0),
        ("  cq_enqueue(): Atomic write to tail → O(1)", 1),
        ("  cq_dequeue(): Atomic read from head → O(1)", 1),
        ("Result: Zero contention between threads", 0)
    ]
)

# Slide 9: Bit Vector (Member 2)
add_content_slide(
    "Bit Vector: Layer 1 Filtering",
    [
        ("Purpose: Fast O(1) IP prefix filtering", 0),
        ("Representation: 64-bit words (efficient for modern CPUs)", 1),
        ("Memory: 2^16 entries = 65,536 IPs = 8 KB", 1),
        ("Operations using bitwise shifts & masks:", 0),
        ("  bv_set(index)     → Set bit to 1 with: words[index >> 6] |= (1 << index & 63)", 1),
        ("  bv_clear(index)   → Clear bit with: words[index >> 6] &= ~(1 << index & 63)", 1),
        ("  bv_contains(index)→ Check bit: (words[index >> 6] >> (index & 63)) & 1", 1),
        ("Use: First-pass filter eliminates ~65K IPs instantly", 0)
    ]
)

# Slide 10: vEB Tree (Member 2)
add_content_slide(
    "vEB Tree: Layer 2 Filtering",
    [
        ("Name: Real-Space van Emde Boas Tree", 0),
        ("Problem: Full 2^32 IP space = 512 MB as flat array", 1),
        ("Solution: Two-tier hierarchical split", 0),
        ("  Top 16 bits: Hash map (clusters selector) - 1024 buckets", 1),
        ("  Bottom 16 bits: BitVector per cluster - 8 KB each", 1),
        ("Memory: On-demand allocation (empty clusters cost 0)", 1),
        ("Operations:", 0),
        ("  insert(ip): Create cluster if needed, set bit → O(1)", 1),
        ("  member(ip): Hash lookup + bit check → O(1)", 1),
        ("Efficiency: 1M IPs occupy ~125 MB instead of 512 MB", 0)
    ]
)

# Slide 11: Gatekeeper (Member 2)
add_content_slide(
    "Member 2: Gatekeeper (IP Filtering)",
    [
        ("Two-Layer Defense Pipeline:", 0),
        ("Layer 1 (Fast): BitVector prefix check", 0),
        ("  - Extract top 16 bits from IP", 1),
        ("  - Check bit in BitVector O(1)", 1),
        ("  - Eliminates ~65,000 IPs per lookup", 1),
        ("Layer 2 (Exact): vEB Tree exact match", 0),
        ("  - Confirms full 32-bit IP address", 1),
        ("  - Avoids false positives from prefix tagging", 1),
        ("Decision:", 0),
        ("  DROP → IP on blacklist or watchlist", 1),
        ("  PASS → IP is safe, forward to scheduler", 1)
    ]
)

# Slide 12: Max Heap Scheduler (Member 3)
add_content_slide(
    "Member 3: Scheduler (Priority Queue)",
    [
        ("Data Structure: Array-based Max Heap", 0),
        ("Capacity: 10,000 packets", 1),
        ("Priority Assignment (1-10 scale):", 0),
        ("  DNS (UDP:53) → Priority 9 (Highest)", 1),
        ("  HTTP/HTTPS (TCP:80/443) → Priority 8", 1),
        ("  Others → Priority 5 (Normal)", 1),
        ("Operations:", 0),
        ("  insert(packet): Add to end + heapify_up → O(log n)", 1),
        ("  extract(): Return max + heapify_down → O(log n)", 1),
        ("Benefit: Critical traffic (DNS, Web) processed first", 0)
    ]
)

# Slide 13: Splay Tree (Member 3 Optional)
add_content_slide(
    "Member 3 (Optional): Splay Tree",
    [
        ("Purpose: Active session tracking and caching", 0),
        ("Self-Balancing BST with splaying optimization", 1),
        ("Time Complexity: O(log n) amortized", 1),
        ("Use Cases:", 0),
        ("  - Quick lookup of recently accessed sessions", 1),
        ("  - Timeout detection and connection expiration", 1),
        ("  - Connection state management", 1),
        ("Status: Design complete, implementation optional for Phase 1", 0),
        ("Future: Can integrate for advanced session tracking", 0)
    ]
)

# Slide 14: Work Completed
add_content_slide(
    "Work Completed",
    [
        ("Core Modules Implemented:", 0),
        ("  [X] Bit Vector - Layer 1 filtering", 1),
        ("  [X] vEB Tree - Layer 2 filtering", 1),
        ("  [X] Gatekeeper - Two-layer IP filter pipeline", 1),
        ("  [X] Concurrent Queue - Lock-free SPSC", 1),
        ("  [X] Capture Engine - libpcap + pthreads", 1),
        ("  [X] Scheduler - Max Heap priority queue", 1),
        ("Testing & Validation:", 0),
        ("  [X] Unit tests - All components", 1),
        ("  [X] Stress test - 1M packets, 100K malicious IPs", 1),
        ("  [X] Capture demo - Real network interface testing", 1),
        ("  [X] Git repository - Tracked on dev branch", 1)
    ]
)

# Slide 15: Work To Be Done
add_content_slide(
    "Work To Be Done",
    [
        ("Core Engine Enhancements:", 0),
        ("  [ ] Main engine integration (tie all modules together)", 1),
        ("  [ ] Splay Tree implementation for session tracking", 1),
        ("  [ ] Fibonacci Heap for advanced QoS scheduling", 1),
        ("ML Integration:", 0),
        ("  [ ] Flow aggregator (convert packets → flows)", 1),
        ("  [ ] Real-time ML model scoring", 1),
        ("  [ ] Distributed multi-threaded flow processing", 1),
        ("Performance & Scale:", 0),
        ("  [ ] Hardware acceleration (XDP/eBPF)", 1),
        ("  [ ] Benchmarking on real traffic (Gbps+)", 1),
        ("  [ ] Production deployment guide", 1)
    ]
)

# Slide 16: Performance Benchmarks
add_content_slide(
    "Performance Benchmarks",
    [
        ("Test Configuration:", 0),
        ("  - 1 Million packets evaluated", 1),
        ("  - 100,000 malicious IPs in watchlist", 1),
        ("Results on MacBook Pro (M1):", 0),
        ("  - Total Time: 0.07 seconds", 1),
        ("  - Throughput: 14.3 million packets/second", 1),
        ("  - Per-packet latency: 70 nanoseconds", 1),
        ("Memory Usage:", 0),
        ("  - BitVector (Layer 1): 8 KB", 1),
        ("  - vEB Tree clusters: ~800 KB (sparse)", 1),
        ("  - Concurrent Queue: 1 KB metadata + buffer", 1)
    ]
)

# Slide 17: Key Innovations
add_content_slide(
    "Key Innovations",
    [
        ("Lock-Free Architecture", 0),
        ("  - Zero mutex overhead, atomic operations only", 1),
        ("  - Sub-microsecond latency vs milliseconds with locks", 1),
        ("Two-Tier vEB Tree", 0),
        ("  - Memory-efficient (on-demand allocation)", 1),
        ("  - Handles full 2^32 IP space efficiently", 1),
        ("Layered IP Filtering", 0),
        ("  - Fast prefix elimination + exact match validation", 1),
        ("  - Avoids false positives", 1),
        ("Real-Time Packet Processing", 0),
        ("  - Millions of packets per second throughput", 1)
    ]
)

# Slide 18: References & Resources
add_content_slide(
    "References & Resources",
    [
        ("Academic Papers:", 0),
        ("  - van Emde Boas, P. (1975) 'Preserving order in a forest'", 1),
        ("  - Herlihy & Shavit (2008) 'The Art of Multiprocessor Programming'", 1),
        ("Libraries & Documentation:", 0),
        ("  - libpcap: https://www.tcpdump.org/", 1),
        ("  - C11 Standard: ISO/IEC 9899:2011", 1),
        ("Project Resources:", 0),
        ("  - GitHub: ArborFlow_ADS_CP", 1),
        ("  - README.md: Architecture details and packet documentation", 1),
        ("  - packet_layer.png: Visual packet layer representation", 1)
    ]
)

# Slide 19: Conclusion
add_title_slide(
    "Thank You!",
    "ArborFlow: Advanced Data Structures in Action\nQuestions?"
)

# Save presentation
output_path = "/Users/prathmesh/ArborFlow_ADS_CP/ArborFlow_Presentation.pptx"
prs.save(output_path)
print(f"✓ Presentation saved to: {output_path}")
print(f"✓ Total slides: 19")
print(f"\nSlide Breakdown:")
print("  1. Title")
print("  2. Project Objectives")
print("  3. Architecture Overview")
print("  4. Tech Stack")
print("  5. Methodology")
print("  6. Literature Review")
print("  7. Packet Collection (Member 1)")
print("  8. Concurrent Queue")
print("  9. Bit Vector (Member 2)")
print(" 10. vEB Tree (Member 2)")
print(" 11. Gatekeeper (Member 2)")
print(" 12. Max Heap Scheduler (Member 3)")
print(" 13. Splay Tree (Member 3 Optional)")
print(" 14. Work Completed")
print(" 15. Work To Be Done")
print(" 16. Performance Benchmarks")
print(" 17. Key Innovations")
print(" 18. References & Resources")
print(" 19. Conclusion")
